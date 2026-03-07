import os
import io
import fitz  # PyMuPDF
from flask import Flask, request, send_file, send_from_directory
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from flask_cors import CORS

app = Flask(__name__, static_folder='.')
CORS(app)

def srgb_to_rgb(srgb):
    """Converts PDF integer color to (R, G, B) tuple."""
    if srgb is None:
        return (0, 0, 0)
    # Extract RGB from the integer color value
    return (srgb >> 16 & 255, srgb >> 8 & 255, srgb & 255)

@app.route('/')
def index():
    """Serves the index.html file from the root directory."""
    return send_from_directory('.', 'index.html')

@app.route("/convert", methods=["POST"])
def convert():
    file = request.files.get("file")
    if not file:
        return "No file uploaded", 400

    # Load PDF from memory
    pdf_bytes = file.read()
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    prs = Presentation()

    for page in doc:
        # --- DYNAMIC LAYOUT LOGIC ---
        # Get PDF page dimensions (72 points = 1 inch)
        pdf_w = page.rect.width
        pdf_h = page.rect.height
        
        # Adjust the PowerPoint slide size to match this specific page
        prs.slide_width = Inches(pdf_w / 72)
        prs.slide_height = Inches(pdf_h / 72)
        
        # Add a blank slide (layout 6 is 'Blank')
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 1. Extract TEXT with exact coordinates, size, and color
        dict_data = page.get_text("dict")
        for block in dict_data["blocks"]:
            if "lines" in block:
                for line in block["lines"]:
                    for span in line["spans"]:
                        l, t, r, b = span["bbox"]
                        width, height = r - l, b - t
                        
                        # Create a textbox at the exact PDF coordinates
                        txt_box = slide.shapes.add_textbox(
                            Inches(l/72), Inches(t/72), 
                            Inches(width/72), Inches(height/72)
                        )
                        tf = txt_box.text_frame
                        tf.word_wrap = True
                        
                        p = tf.paragraphs[0]
                        p.text = span["text"]
                        
                        # Apply Font Size and Color
                        font = p.font
                        font.size = Pt(span["size"])
                        r_val, g_val, b_val = srgb_to_rgb(span["color"])
                        font.color.rgb = RGBColor(r_val, g_val, b_val)

        # 2. Extract IMAGES as standalone movable objects
        for img in page.get_images(full=True):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_stream = io.BytesIO(base_image["image"])
            
            # Find where this image is placed on the page
            img_rects = page.get_image_rects(xref)
            if img_rects:
                rect = img_rects[0]
                slide.shapes.add_picture(
                    image_stream, 
                    Inches(rect.x0/72), 
                    Inches(rect.y0/72), 
                    width=Inches(rect.width/72),
                    height=Inches(rect.height/72)
                )

    # Save the PowerPoint to an in-memory buffer
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)

    return send_file(
        output, 
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation", 
        as_attachment=True, 
        download_name="Converted_Presentation.pptx"
    )

if __name__ == "__main__":
    # Render requires the app to listen on a specific environment PORT
    port = int(os.environ.get("PORT", 5000))
    app.run(host="0.0.0.0", port=port)
